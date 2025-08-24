from pptx import Presentation
import base64
from io import BytesIO

def read_pptx(data_base64):
    try:
        ppt_bytes = base64.urlsafe_b64decode(data_base64)
        ppt_file = BytesIO(ppt_bytes)
        prs = Presentation(ppt_file)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text.strip()
    except Exception as e:
        return f"[Error reading PPTX] {e}"
